#!./.env/bin/python
# -*- coding: utf-8 -*-
# Time-stamp: "2024-01-14 12:56:44 (ywatanabe)"

# pptx2tif.py

# pip install aspose.slides

import sys
from pptx import Presentation
from PIL import Image
import io
import cv2
import numpy as np
import os

#!./.env/bin/python
# -*- coding: utf-8 -*-

# pptx2tif.py

import sys
from pptx import Presentation
from PIL import Image
import io
import cv2
import numpy as np
import os

def find_content_area(image_path):  # [REVISED]
    img = cv2.imread(image_path)
    if img is None:
        raise FileNotFoundError(f"Unable to read the image file: {image_path}")
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY_INV)
    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    max_contour = max(contours, key=cv2.contourArea)
    x, y, w, h = cv2.boundingRect(max_contour)
    return x, y, x + w, y + h

def convert_pptx_to_tif(pptx_path):  # [REVISED]
    import ipdb; ipdb.set_trace()
    pres = Presentation(pptx_path)
    for i, slide in enumerate(pres.slides):
        print()
        # Export each slide as an image
        image = Image.open(io.BytesIO(slide.shapes[0].image.blob))
        # Save the image as TIFF
        image.save(f"slide_{i}.tif", format='TIFF')  # [REVISED]



# def find_content_area(image_path):
#     img = cv2.imread(image_path)
#     if img is None:
#         raise FileNotFoundError(f"Unable to read the image file: {image_path}")
#     gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
#     _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY_INV)
#     contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
#     max_contour = max(contours, key=cv2.contourArea)
#     x, y, w, h = cv2.boundingRect(max_contour)
#     return x, y, x + w, y + h

# def convert_pptx_to_tif(pptx_path):
#     pres = Presentation(pptx_path)
#     first_slide = pres.slides[0]
#     slide_image_stream = io.BytesIO()
#     # Use a renderer to save the slide as an image
#     # This requires a rendering library or tool that can convert slides to images
#     # Placeholder for actual slide rendering code
#     # first_slide.get_image().save(slide_image_stream, format='TIFF')

#     # Assuming the rendering tool saves the slide as 'slide_image.png'
#     temp_image_path = "slide_image.png"

#     # Render the first slide to an image file here
#     # This is a placeholder for the rendering code
#     # You need to implement the rendering of the slide to an image file

#     crop_area = find_content_area(temp_image_path)
#     image = Image.open(temp_image_path)
#     cropped_image = image.crop(crop_area)

#     base_name = os.path.splitext(os.path.basename(pptx_path))[0]
#     output_path = f"{base_name}.tif"
#     cropped_image.save(output_path, format='TIFF')
#     os.remove(temp_image_path)

if __name__ == "__main__":
    # sudo apt-get install unoconv -y
    # sudo apt-get install libreoffice -y
    # sudo apt-get install libreoffice-dev -y        
    unoconv -f tiff ./Figure_01.pptx

    
    pptx_path = "./src/figures/Figure_01.pptx"
    # convert_pptx_to_tif(pptx_path)
    
    import aspose.slides as slides
    presentation = slides.Presentation(pptx_path)
    presentation.save(pptx_path.replace(".pptx", ".tif"), slides.export.SaveFormat.TIFF)
    # RuntimeError: Proxy error(TypeInitializationException): The type initializer for 'Gdip' threw an exception. ---> DllNotFoundException: Unable to load shared library 'libgdiplus' or one of its dependencies. In order to help diagnose loading problems, consider setting the LD_DEBUG environment variable: liblibgdiplus: cannot open shared object file: No such file or directory
    sudo apt-get install -y libgdiplus

    # # pptx_path = sys.argv[1]
    # convert_pptx_to_tif(pptx_path)
    # # ./.env/bin/python ./.scripts/py/pptx2tif.py ./src/figures/Figure_01.pptx
